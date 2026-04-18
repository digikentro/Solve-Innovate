from __future__ import annotations

import asyncio
import base64
import uuid
from pathlib import Path
from typing import AsyncGenerator

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine
from sqlmodel import SQLModel

from api.v1.ppt.endpoints.markdown_presentation import MARKDOWN_ROUTER
from api.v1.ppt.endpoints.spatial_presentation import SPATIAL_ROUTER
from models.sql.markdown_presentation import MarkdownPresentationModel
from models.sql.project_presentation import ProjectPresentationModel
from models.sql.project_presentation_revision import ProjectPresentationRevisionModel
from services.database import get_async_session
from services.spatial_pptx_export import (
    _extract_text_runs,
    _normalize_position,
    _percent_to_emu,
    _render_image_block,
)


@pytest.fixture
def test_app(async_session_factory: async_sessionmaker[AsyncSession]) -> FastAPI:
    app = FastAPI()
    app.include_router(MARKDOWN_ROUTER, prefix="/api/v1/ppt")
    app.include_router(SPATIAL_ROUTER, prefix="/api/v1/ppt")

    async def _override_get_async_session() -> AsyncGenerator[AsyncSession, None]:
        async with async_session_factory() as session:
            yield session

    app.dependency_overrides[get_async_session] = _override_get_async_session
    return app


@pytest.fixture
def client(test_app: FastAPI) -> TestClient:
    return TestClient(test_app)


@pytest.fixture
def async_session_factory(tmp_path: Path) -> async_sessionmaker[AsyncSession]:
    db_path = tmp_path / "spatial_export_tests.db"
    engine = create_async_engine(f"sqlite+aiosqlite:///{db_path}")

    async def _create_tables() -> None:
        async with engine.begin() as conn:
            await conn.run_sync(
                lambda sync_conn: SQLModel.metadata.create_all(
                    sync_conn,
                    tables=[
                        MarkdownPresentationModel.__table__,
                        ProjectPresentationModel.__table__,
                        ProjectPresentationRevisionModel.__table__,
                    ],
                )
            )

    asyncio.run(_create_tables())

    session_factory = async_sessionmaker(engine, expire_on_commit=False)
    try:
        yield session_factory
    finally:
        asyncio.run(engine.dispose())


@pytest.fixture
def seeded_ids(async_session_factory: async_sessionmaker[AsyncSession]) -> dict[str, str]:
    markdown_id = uuid.uuid4()
    project_id = uuid.uuid4()
    revision_id = uuid.uuid4()

    editor_payload = {
        "format": "spatial-json-canvas",
        "version": "1.0",
        "slides": [
            {
                "id": "slide-1",
                "title": "Test Slide",
                "visual_intent": "Narrative",
                "blocks": [
                    {
                        "id": "text-1",
                        "type": "text",
                        "position": {"x": 10, "y": 10, "width": 50, "height": 20},
                        "content": "<p>Hello world</p>",
                        "style": {"font_size": 20, "align": "left"},
                    }
                ],
            }
        ],
    }

    async def _seed() -> None:
        async with async_session_factory() as session:
            markdown = MarkdownPresentationModel(
                id=markdown_id,
                content="seed content",
                n_slides=1,
                language="English",
                theme="modern",
            )
            project = ProjectPresentationModel(
                id=project_id,
                project_id="project-seed",
                title="Spatial Export Test",
                status="draft",
                current_revision_id=revision_id,
            )
            revision = ProjectPresentationRevisionModel(
                id=revision_id,
                presentation_id=project_id,
                revision_number=1,
                is_draft=True,
                markdown_presentation_id=markdown_id,
                settings={"editor_payload": editor_payload},
            )

            session.add(markdown)
            session.add(project)
            session.add(revision)
            await session.commit()

    asyncio.run(_seed())

    return {
        "markdown_id": str(markdown_id),
        "project_id": str(project_id),
    }


def test_coordinate_math_middle_to_bottom_right_quadrant() -> None:
    prs = Presentation()
    left, top, width, height = _percent_to_emu(prs, 50.0, 50.0, 50.0, 50.0)

    assert int(left) == int(Inches(5.0))
    # Default built-in template is 10" × 7.5" (4:3), not 16:9.
    assert int(top) == int(Inches(3.75))
    assert int(width) == int(Inches(5.0))
    assert int(height) == int(Inches(3.75))


def test_missing_coordinates_fallback_without_keyerror() -> None:
    x, y, w, h = _normalize_position({"x": None, "y": None, "width": 40, "height": 20})

    # Missing x/y falls back to centered coordinates for provided width/height.
    assert x == pytest.approx(30.0)
    assert y == pytest.approx(40.0)
    assert w == pytest.approx(40.0)
    assert h == pytest.approx(20.0)


def test_html_stripping_and_text_run_translation() -> None:
    lines = _extract_text_runs("<p><strong>Bold</strong> text</p>")

    assert len(lines) >= 1
    first_line = lines[0]
    reconstructed = "".join(run[0] for run in first_line)
    assert "Bold" in reconstructed
    assert "text" in reconstructed
    assert "<" not in reconstructed
    assert any(run[0].strip() == "Bold" and run[1] is True for run in first_line)


def test_missing_images_logs_warning_and_continues(caplog: pytest.LogCaptureFixture) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    before_shape_count = len(slide.shapes)

    broken_image_block = {
        "id": "img-broken",
        "type": "image",
        "prompt": "C:/definitely/not/real/nonexistent_image.png",
    }

    with caplog.at_level("WARNING"):
        asyncio.run(
            _render_image_block(
                slide,
                broken_image_block,
                left=Inches(1),
                top=Inches(1),
                width=Inches(2),
                height=Inches(2),
            )
        )

    assert len(slide.shapes) == before_shape_count
    assert "Skipping image block 'img-broken'" in caplog.text


def test_image_block_embeds_data_uri(monkeypatch: pytest.MonkeyPatch) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    png_data_uri = (
        "data:image/png;base64,"
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jWZkAAAAASUVORK5CYII="
    )

    asyncio.run(
        _render_image_block(
            slide,
            {"id": "img-data", "type": "image", "prompt": png_data_uri},
            left=Inches(1),
            top=Inches(1),
            width=Inches(1),
            height=Inches(1),
        )
    )

    assert len(slide.shapes) > 0


def test_image_block_embeds_remote_url(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_file = tmp_path / "remote.png"
    image_file.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+jWZkAAAAASUVORK5CYII="
        )
    )

    async def _fake_download(url: str, save_directory: str, headers=None):
        return str(image_file)

    monkeypatch.setattr("services.spatial_pptx_export.download_file", _fake_download)

    asyncio.run(
        _render_image_block(
            slide,
            {"id": "img-remote", "type": "image", "prompt": "https://example.com/test.png"},
            left=Inches(1),
            top=Inches(1),
            width=Inches(1),
            height=Inches(1),
        )
    )

    assert len(slide.shapes) > 0


def test_markdown_export_endpoint_returns_file_download(
    client: TestClient,
    seeded_ids: dict[str, str],
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    expected_path = str(tmp_path / f"spatial_presentation_{seeded_ids['markdown_id']}.pptx")

    async def _fake_export(*, editor_payload, output_path: str) -> str:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"PPTX")
        return output_path

    monkeypatch.setattr(
        "services.spatial_pptx_export.export_spatial_to_pptx",
        _fake_export,
    )
    monkeypatch.setattr(
        "api.v1.ppt.endpoints.markdown_presentation.get_exports_directory",
        lambda: str(tmp_path),
    )

    response = client.post(
        f"/api/v1/ppt/markdown/presentation/{seeded_ids['markdown_id']}/export/pptx"
    )

    assert response.status_code == 200
    assert response.content == b"PPTX"
    cd = response.headers.get("content-disposition") or ""
    assert "attachment" in cd.lower()
    assert Path(expected_path).exists()


def test_spatial_export_endpoint_returns_200_and_path(
    client: TestClient,
    seeded_ids: dict[str, str],
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    async def _fake_export(*, editor_payload, output_path: str) -> str:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"PPTX")
        return output_path

    monkeypatch.setattr(
        "api.v1.ppt.endpoints.spatial_presentation.export_spatial_to_pptx",
        _fake_export,
    )
    monkeypatch.setattr(
        "api.v1.ppt.endpoints.spatial_presentation.get_exports_directory",
        lambda: str(tmp_path),
    )

    response = client.post(
        f"/api/v1/ppt/spatial/presentation/{seeded_ids['project_id']}/export/pptx"
    )

    assert response.status_code == 200
    body = response.json()
    assert isinstance(body.get("path"), str)
    assert body["path"].endswith(".pptx")
    assert Path(body["path"]).exists()
