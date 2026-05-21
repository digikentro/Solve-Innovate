"""
markdown_pptx_export.py

Converts a list of extended-markdown slide strings to a .pptx file.

Design notes
────────────
• Each slide string is produced by the markdown generation pipeline.
• Parsing is intentionally simple and line-based — the same extended syntax
  used in the frontend renderer is handled here.
• This service intentionally has no FastAPI dependency so it can be unit-tested
  in isolation.
"""

from __future__ import annotations

import re
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Theme definitions ────────────────────────────────────────────────────────

THEMES: Dict[str, Dict[str, Any]] = {
    "modern-dark": {
        "bg": "1a1a2e",
        "primary": "5146E5",
        "text": "ffffff",
        "subtext": "b0b0c8",
    },
    "clean-light": {
        "bg": "ffffff",
        "primary": "2563eb",
        "text": "1f2937",
        "subtext": "6b7280",
    },
    "nature": {
        "bg": "ecfdf5",
        "primary": "059669",
        "text": "064e3b",
        "subtext": "065f46",
    },
    "sunset": {
        "bg": "fffbeb",
        "primary": "f59e0b",
        "text": "78350f",
        "subtext": "92400e",
    },
    "midnight": {
        "bg": "0f0f1a",
        "primary": "8b5cf6",
        "text": "e2e8f0",
        "subtext": "94a3b8",
    },
}

DEFAULT_THEME = THEMES["modern-dark"]

# ─── Slide dimensions (16:9) ──────────────────────────────────────────────────

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)
MARGIN = Inches(0.75)
CONTENT_W = Inches(14.5)


# ─── Colour helpers ───────────────────────────────────────────────────────────


def _rgb(hex_str: str) -> RGBColor:
    """Convert 6-char hex string (no #) to RGBColor."""
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# ─── Inline markdown → plain text ────────────────────────────────────────────

_INLINE_RE = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`")


def _plain(text: str) -> str:
    """Strip inline bold / italic / code markers from text."""
    return _INLINE_RE.sub(lambda m: m.group(1) or m.group(2) or m.group(3), text)


# ─── Markdown parser ──────────────────────────────────────────────────────────


def _parse_slide(markdown: str) -> List[Dict[str, Any]]:
    """
    Parse an extended-markdown slide string into a list of block dicts.

    Supported block types:
      heading    – { type, level, text }
      paragraph  – { type, text }
      bullet     – { type, items: [str] }
      numbered   – { type, items: [str] }
      quote      – { type, text, attribution? }
      metric     – { type, metrics: [{value, label, delta?}] }
      callout    – { type, icon, title, body }
      chart      – { type, chart_type, headers, rows }
      table      – { type, headers, rows }
      image      – { type, alt, src }
      columns    – { type, columns: [[blocks]] }  (simplified — treated as paragraph)
    """
    blocks: List[Dict[str, Any]] = []
    lines = markdown.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i]

        # ── Headings ──────────────────────────────────────────────────────────
        if line.startswith("### "):
            blocks.append({"type": "heading", "level": 3, "text": _plain(line[4:])})
            i += 1
            continue
        if line.startswith("## "):
            blocks.append({"type": "heading", "level": 2, "text": _plain(line[3:])})
            i += 1
            continue
        if line.startswith("# "):
            blocks.append({"type": "heading", "level": 1, "text": _plain(line[2:])})
            i += 1
            continue

        # ── Metric block ──────────────────────────────────────────────────────
        if line.strip() == ":::metric":
            metrics: List[Dict[str, str]] = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith(":::"):
                parts = [p.strip() for p in lines[i].split("|")]
                if len(parts) >= 2:
                    m: Dict[str, str] = {"value": parts[0], "label": parts[1]}
                    if len(parts) >= 3 and parts[2]:
                        m["delta"] = parts[2]
                    metrics.append(m)
                i += 1
            if metrics:
                blocks.append({"type": "metric", "metrics": metrics})
            i += 1  # skip closing :::
            continue

        # ── Callout block ─────────────────────────────────────────────────────
        callout_m = re.match(r"^:::callout\[icon=(\w+)\]", line.strip())
        if callout_m:
            icon = callout_m.group(1)
            i += 1
            title = ""
            body_lines: List[str] = []
            while i < len(lines) and not lines[i].strip().startswith(":::"):
                l = lines[i]
                bold_m = re.match(r"^\*\*(.+)\*\*$", l.strip())
                if bold_m:
                    title = bold_m.group(1)
                else:
                    body_lines.append(_plain(l))
                i += 1
            blocks.append(
                {
                    "type": "callout",
                    "icon": icon,
                    "title": title,
                    "body": "\n".join(body_lines).strip(),
                }
            )
            i += 1  # skip closing :::
            continue

        # ── Chart block ───────────────────────────────────────────────────────
        chart_m = re.match(r"^:::chart\[type=(\w+)\]", line.strip())
        if chart_m:
            chart_type = chart_m.group(1)
            i += 1
            headers: List[str] = []
            rows: List[List[str]] = []
            while i < len(lines) and not lines[i].strip().startswith(":::"):
                row_line = lines[i]
                if "|" in row_line:
                    cells = [c.strip() for c in row_line.split("|") if c.strip()]
                    if not any(re.match(r"^-+$", c) for c in cells):
                        if not headers:
                            headers = cells
                        else:
                            rows.append(cells)
                i += 1
            blocks.append(
                {
                    "type": "chart",
                    "chart_type": chart_type,
                    "headers": headers,
                    "rows": rows,
                }
            )
            i += 1
            continue

        # ── Columns block ─────────────────────────────────────────────────────
        if line.strip() == ":::columns":
            col_text: List[str] = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith(":::"):
                col_text.append(lines[i])
                i += 1
            # Flatten to a paragraph for PPTX (columns not supported natively)
            blocks.append({"type": "paragraph", "text": _plain(" ".join(col_text))})
            i += 1
            continue

        # ── Block quotes ──────────────────────────────────────────────────────
        if line.startswith("> "):
            text_parts: List[str] = []
            attribution: Optional[str] = None
            while i < len(lines) and lines[i].startswith("> "):
                content = lines[i][2:]
                attr_m = re.match(r"^[—–-]\s*(.+)$", content)
                if attr_m:
                    attribution = attr_m.group(1)
                else:
                    text_parts.append(_plain(content.strip('"').strip("\u201c\u201d")))
                i += 1
            blocks.append(
                {
                    "type": "quote",
                    "text": " ".join(text_parts),
                    "attribution": attribution,
                }
            )
            continue

        # ── Bullet list ───────────────────────────────────────────────────────
        if line.startswith("- "):
            items: List[str] = []
            while i < len(lines) and lines[i].startswith("- "):
                items.append(_plain(lines[i][2:]))
                i += 1
            blocks.append({"type": "bullet", "items": items})
            continue

        # ── Numbered list ─────────────────────────────────────────────────────
        if re.match(r"^\d+\.\s", line):
            items = []
            while i < len(lines) and re.match(r"^\d+\.\s", lines[i]):
                items.append(_plain(re.sub(r"^\d+\.\s*", "", lines[i])))
                i += 1
            blocks.append({"type": "numbered", "items": items})
            continue

        # ── Table ─────────────────────────────────────────────────────────────
        if "|" in line and i + 1 < len(lines) and "---" in lines[i + 1]:
            headers = [c.strip() for c in line.split("|") if c.strip()]
            i += 2  # skip separator
            rows = []
            while i < len(lines) and "|" in lines[i]:
                rows.append([c.strip() for c in lines[i].split("|") if c.strip()])
                i += 1
            blocks.append({"type": "table", "headers": headers, "rows": rows})
            continue

        # ── Image ─────────────────────────────────────────────────────────────
        img_m = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
        if img_m:
            blocks.append({"type": "image", "alt": img_m.group(1), "src": img_m.group(2)})
            i += 1
            continue

        # ── Paragraph ─────────────────────────────────────────────────────────
        if line.strip():
            blocks.append({"type": "paragraph", "text": _plain(line.strip())})
        i += 1

    return blocks


# ─── PPTX rendering helpers ───────────────────────────────────────────────────


def _set_bg(slide, hex_color: str) -> None:
    """Fill slide background with a solid colour."""
    from pptx.oxml.ns import qn
    from lxml import etree  # bundled with python-pptx

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Any:
    return slide.shapes.add_textbox(left, top, width, height)


def _para(tf, text: str, font_size: int, bold: bool, color: RGBColor, align=PP_ALIGN.LEFT):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return p


def _render_slide(prs: Presentation, slide_md: str, theme: Dict[str, Any]) -> None:
    """Add a single slide to the Presentation object."""
    bg_col = theme.get("bg", "1a1a2e")
    primary_col = _rgb(theme.get("primary", "5146E5"))
    text_col = _rgb(theme.get("text", "ffffff"))
    subtext_col = _rgb(theme.get("subtext", "b0b0c8"))
    positive_col = _rgb("10b981")
    negative_col = _rgb("ef4444")

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, bg_col)

    blocks = _parse_slide(slide_md)
    y = MARGIN

    for block in blocks:
        btype = block["type"]

        # ── Heading ──────────────────────────────────────────────────────────
        if btype == "heading":
            level = block["level"]
            font_size = {1: 44, 2: 32, 3: 24}.get(level, 24)
            height = Inches(0.8 if level > 1 else 1.1)
            tb = _add_textbox(slide, MARGIN, y, CONTENT_W, height)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = block["text"]
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = primary_col
            y += height + Inches(0.05)

        # ── Paragraph ────────────────────────────────────────────────────────
        elif btype == "paragraph":
            height = Inches(0.5)
            tb = _add_textbox(slide, MARGIN, y, CONTENT_W, height)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = block["text"]
            run.font.size = Pt(20)
            run.font.color.rgb = text_col
            y += height + Inches(0.05)

        # ── Bullet list ───────────────────────────────────────────────────────
        elif btype == "bullet":
            items: List[str] = block["items"]
            height = Inches(0.4 * len(items) + 0.1)
            tb = _add_textbox(slide, MARGIN + Inches(0.2), y, CONTENT_W - Inches(0.2), height)
            tf = tb.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(items):
                p = tf.paragraphs[0] if idx == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
                run = p.add_run()
                run.text = f"• {item}"
                run.font.size = Pt(20)
                run.font.color.rgb = text_col
            y += height + Inches(0.05)

        # ── Numbered list ─────────────────────────────────────────────────────
        elif btype == "numbered":
            items = block["items"]
            height = Inches(0.4 * len(items) + 0.1)
            tb = _add_textbox(slide, MARGIN + Inches(0.2), y, CONTENT_W - Inches(0.2), height)
            tf = tb.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(items):
                p = tf.paragraphs[0] if idx == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
                run = p.add_run()
                run.text = f"{idx + 1}. {item}"
                run.font.size = Pt(20)
                run.font.color.rgb = text_col
            y += height + Inches(0.05)

        # ── Quote ─────────────────────────────────────────────────────────────
        elif btype == "quote":
            height = Inches(1.5)
            tb = _add_textbox(slide, MARGIN + Inches(0.4), y, CONTENT_W - Inches(0.4), height)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f'"{block["text"]}"'
            run.font.size = Pt(26)
            run.font.italic = True
            run.font.color.rgb = text_col
            if block.get("attribution"):
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = f"— {block['attribution']}"
                r2.font.size = Pt(18)
                r2.font.color.rgb = subtext_col
            y += height + Inches(0.1)

        # ── Metric row ────────────────────────────────────────────────────────
        elif btype == "metric":
            metrics: List[Dict[str, str]] = block["metrics"]
            n = len(metrics)
            box_w = CONTENT_W / n
            row_h = Inches(1.6)
            for idx, metric in enumerate(metrics):
                x = MARGIN + idx * box_w
                tb = _add_textbox(slide, x, y, box_w - Inches(0.1), row_h)
                tf = tb.text_frame
                tf.word_wrap = True
                # Value
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = metric["value"]
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = primary_col
                p.alignment = PP_ALIGN.CENTER
                # Label
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = metric["label"]
                r2.font.size = Pt(18)
                r2.font.color.rgb = subtext_col
                p2.alignment = PP_ALIGN.CENTER
                # Delta
                if metric.get("delta"):
                    p3 = tf.add_paragraph()
                    r3 = p3.add_run()
                    r3.text = metric["delta"]
                    r3.font.size = Pt(16)
                    r3.font.color.rgb = (
                        positive_col if metric["delta"].startswith("+") else negative_col
                    )
                    p3.alignment = PP_ALIGN.CENTER
            y += row_h + Inches(0.15)

        # ── Callout ───────────────────────────────────────────────────────────
        elif btype == "callout":
            height = Inches(1.2)
            tb = _add_textbox(slide, MARGIN, y, CONTENT_W, height)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            icon_map = {
                "lightbulb": "💡",
                "warning": "⚠️",
                "info": "ℹ️",
                "check": "✅",
            }
            icon_str = icon_map.get(block.get("icon", ""), "•")
            if block.get("title"):
                run = p.add_run()
                run.text = f"{icon_str} {block['title']}"
                run.font.size = Pt(22)
                run.font.bold = True
                run.font.color.rgb = primary_col
            if block.get("body"):
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = block["body"]
                r2.font.size = Pt(18)
                r2.font.color.rgb = text_col
            y += height + Inches(0.1)

        # ── Chart (rendered as a data table in PPTX) ──────────────────────────
        elif btype == "chart":
            headers: List[str] = block.get("headers", [])
            rows: List[List[str]] = block.get("rows", [])
            if headers:
                # Render as a table shape
                row_count = len(rows) + 1  # +1 for header
                col_count = len(headers)
                if row_count > 1 and col_count > 0:
                    col_w = CONTENT_W / col_count
                    row_h = Inches(0.4)
                    table_h = row_h * row_count
                    tbl = slide.shapes.add_table(
                        row_count, col_count,
                        int(MARGIN), int(y),
                        int(CONTENT_W), int(table_h)
                    ).table
                    # Header row
                    for ci, hdr in enumerate(headers):
                        cell = tbl.cell(0, ci)
                        cell.text = hdr
                        cell.text_frame.paragraphs[0].runs[0].font.bold = True
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
                        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = primary_col
                    # Data rows
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row[:col_count]):
                            cell = tbl.cell(ri + 1, ci)
                            cell.text = val
                            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = text_col
                    y += table_h + Inches(0.15)

        # ── Table ─────────────────────────────────────────────────────────────
        elif btype == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            if headers:
                row_count = len(rows) + 1
                col_count = len(headers)
                if row_count > 1 and col_count > 0:
                    row_h = Inches(0.4)
                    table_h = row_h * row_count
                    tbl = slide.shapes.add_table(
                        row_count, col_count,
                        int(MARGIN), int(y),
                        int(CONTENT_W), int(table_h)
                    ).table
                    for ci, hdr in enumerate(headers):
                        cell = tbl.cell(0, ci)
                        cell.text = hdr
                        try:
                            cell.text_frame.paragraphs[0].runs[0].font.bold = True
                            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
                            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = primary_col
                        except IndexError:
                            pass
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row[:col_count]):
                            cell = tbl.cell(ri + 1, ci)
                            cell.text = val
                            try:
                                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = text_col
                            except IndexError:
                                pass
                    y += table_h + Inches(0.15)

        # ── Image (placeholder text — actual images need URL fetching) ────────
        elif btype == "image":
            alt = block.get("alt", "")
            src = block.get("src", "")
            # If it's an AI image placeholder, just note it as text
            if src.startswith("image:"):
                tb = _add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"[Image: {src[6:] or alt}]"
                run.font.size = Pt(16)
                run.font.italic = True
                run.font.color.rgb = subtext_col
                y += Inches(0.5)

        # Any unknown block type — skip silently

    # Guard: make sure we don't overflow the slide silently
    _ = y  # y is not used after loop — kept for layout tracking


# ─── Public API ───────────────────────────────────────────────────────────────


async def export_markdown_to_pptx(
    slides_markdown: List[str],
    theme_id: str,
    output_path: str,
) -> str:
    """
    Convert a list of extended-markdown slide strings to a .pptx file.

    Parameters
    ----------
    slides_markdown : One string per slide (already split on ---).
    theme_id        : One of the keys in THEMES (falls back to modern-dark).
    output_path     : Absolute path for the output file.

    Returns
    -------
    The absolute path to the saved .pptx file.
    """
    theme = THEMES.get(theme_id, DEFAULT_THEME)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_md in slides_markdown:
        _render_slide(prs, slide_md.strip(), theme)

    import os  # noqa: PLC0415

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
