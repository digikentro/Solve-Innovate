from __future__ import annotations

import base64
import logging
import os
import re
import tempfile
from html import unescape
from html.parser import HTMLParser
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import unquote

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from utils.download_helpers import download_file
from utils.get_env import get_app_data_directory_env

logger = logging.getLogger(__name__)

# Minimum shape size (~0.01 in); avoids invalid zero-area OOXML that PowerPoint rejects.
_MIN_SHAPE_EMU = 9525


def _clamp(value: float, minimum: float, maximum: float) -> float:
    return max(minimum, min(maximum, value))


def _as_float(value: Any, fallback: float) -> float:
    try:
        parsed = float(value)
        if parsed != parsed:  # NaN
            return fallback
        return parsed
    except (TypeError, ValueError):
        return fallback


def _normalize_position(position: Any) -> Tuple[float, float, float, float]:
    """Return normalized (x, y, w, h) percentages on [0, 100] with defensive defaults.

    If missing/invalid, defaults to centered box.
    """
    if not isinstance(position, dict):
        x = 10.0
        y = 35.0
        w = 80.0
        h = 30.0
        return x, y, w, h

    w = _clamp(_as_float(position.get("width"), 80.0), 1.0, 100.0)
    h = _clamp(_as_float(position.get("height"), 30.0), 1.0, 100.0)
    x = _clamp(_as_float(position.get("x"), (100.0 - w) / 2.0), 0.0, 100.0 - w)
    y = _clamp(_as_float(position.get("y"), (100.0 - h) / 2.0), 0.0, 100.0 - h)
    return x, y, w, h


def _percent_to_emu(
    prs: Presentation, x: float, y: float, w: float, h: float
) -> Tuple[Emu, Emu, Emu, Emu]:
    """Map block percentages to EMU using *this* presentation's slide size.

    Do not change ``prs.slide_width`` / ``slide_height`` to a different aspect ratio
    without updating masters: that mismatch commonly triggers PowerPoint "repair"
    dialogs and failed opens. The web editor uses percentage layout; scaling to the
    built-in template's dimensions keeps the package consistent.
    """
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    left = int(round(x / 100.0 * sw))
    top = int(round(y / 100.0 * sh))
    width = max(_MIN_SHAPE_EMU, int(round(w / 100.0 * sw)))
    height = max(_MIN_SHAPE_EMU, int(round(h / 100.0 * sh)))
    left = max(0, min(left, sw - width))
    top = max(0, min(top, sh - height))
    return Emu(left), Emu(top), Emu(width), Emu(height)


def _blank_slide_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[6]


def _parse_hex_color(value: Any) -> Optional[RGBColor]:
    if not isinstance(value, str):
        return None
    candidate = value.strip()
    if candidate.startswith("#"):
        candidate = candidate[1:]
    if len(candidate) == 3:
        candidate = "".join(ch * 2 for ch in candidate)
    if not re.fullmatch(r"[0-9a-fA-F]{6}", candidate):
        return None
    return RGBColor(int(candidate[0:2], 16), int(candidate[2:4], 16), int(candidate[4:6], 16))


class _TipTapHTMLParser(HTMLParser):
    """Minimal TipTap HTML parser preserving paragraph breaks and <b>/<i> runs."""

    def __init__(self) -> None:
        super().__init__()
        self._bold_depth = 0
        self._italic_depth = 0
        self.tokens: List[Tuple[str, bool, bool]] = []

    def handle_starttag(self, tag: str, attrs: List[Tuple[str, Optional[str]]]) -> None:
        tag = tag.lower()
        if tag in {"strong", "b"}:
            self._bold_depth += 1
        elif tag in {"em", "i"}:
            self._italic_depth += 1
        elif tag in {"br", "p", "div", "li"}:
            self.tokens.append(("\n", False, False))

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"strong", "b"}:
            self._bold_depth = max(0, self._bold_depth - 1)
        elif tag in {"em", "i"}:
            self._italic_depth = max(0, self._italic_depth - 1)
        elif tag in {"p", "div", "li"}:
            self.tokens.append(("\n", False, False))

    def handle_data(self, data: str) -> None:
        text = unescape(data)
        if not text:
            return
        self.tokens.append((text, self._bold_depth > 0, self._italic_depth > 0))


def _extract_text_runs(content: Any) -> List[List[Tuple[str, bool, bool]]]:
    """Return lines of (text, bold, italic) runs from TipTap HTML."""
    if not isinstance(content, str):
        return [[("", False, False)]]

    parser = _TipTapHTMLParser()
    try:
        parser.feed(content)
        parser.close()
    except Exception:  # noqa: BLE001
        # Fall back to plain text strip if malformed HTML.
        plain = re.sub(r"<[^>]+>", "", content)
        return [[(plain, False, False)]]

    lines: List[List[Tuple[str, bool, bool]]] = [[]]
    for text, bold, italic in parser.tokens:
        if text == "\n":
            if lines[-1]:
                lines.append([])
            continue
        parts = text.split("\n")
        for idx, part in enumerate(parts):
            if part:
                lines[-1].append((part, bold, italic))
            if idx < len(parts) - 1:
                lines.append([])

    cleaned = [line for line in lines if line]
    return cleaned or [[("", False, False)]]


def _resolve_image_path(raw_path: Any) -> Optional[str]:
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None

    candidate = raw_path.strip()
    if os.path.exists(candidate):
        return candidate

    # Handle frontend/static style references.
    app_data_dir = get_app_data_directory_env() or ""
    if candidate.startswith("/app_data/") and app_data_dir:
        normalized_rel = candidate[len("/app_data/") :].replace("/", os.sep)
        computed = os.path.join(app_data_dir, normalized_rel)
        if os.path.exists(computed):
            return computed

    # If only basename was persisted, try app_data/images.
    if app_data_dir:
        base_name = os.path.basename(candidate)
        images_guess = os.path.join(app_data_dir, "images", base_name)
        if os.path.exists(images_guess):
            return images_guess

    return None


async def _materialize_image_source(raw_path: Any) -> Optional[str]:
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None

    candidate = raw_path.strip()
    resolved_path = _resolve_image_path(candidate)
    if resolved_path:
        return resolved_path

    if candidate.startswith("data:image/"):
        header, _, encoded = candidate.partition(",")
        if not encoded:
            return None

        suffix = ".png"
        match = re.match(r"data:image/([a-zA-Z0-9.+-]+)", header)
        if match:
            mime = match.group(1).lower()
            if mime in {"jpeg", "jpg"}:
                suffix = ".jpg"
            elif mime == "gif":
                suffix = ".gif"
            elif mime == "webp":
                suffix = ".webp"
            elif mime == "svg+xml":
                suffix = ".svg"

        payload = base64.b64decode(encoded) if ";base64" in header else unquote(encoded).encode("utf-8")
        temp_dir = os.path.join(tempfile.gettempdir(), "solve_innovate_pptx_images")
        os.makedirs(temp_dir, exist_ok=True)
        file_path = os.path.join(temp_dir, f"image_{abs(hash(candidate))}{suffix}")
        with open(file_path, "wb") as file_handle:
            file_handle.write(payload)
        return file_path

    if candidate.startswith(("http://", "https://")):
        temp_dir = os.path.join(tempfile.gettempdir(), "solve_innovate_pptx_images")
        downloaded = await download_file(candidate, temp_dir)
        if downloaded:
            return downloaded

    return None


def _render_text_block(slide, block: Dict[str, Any], left: Emu, top: Emu, width: Emu, height: Emu) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    style = block.get("style") if isinstance(block.get("style"), dict) else {}
    color = _parse_hex_color(style.get("color"))
    font_size = _as_float(style.get("font_size"), _as_float(style.get("fontSize"), 18.0))
    align_raw = str(style.get("align") or "left").lower()
    alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align_raw, PP_ALIGN.LEFT)

    lines = _extract_text_runs(block.get("content", ""))
    for line_index, runs in enumerate(lines):
        paragraph = tf.paragraphs[0] if line_index == 0 else tf.add_paragraph()
        paragraph.alignment = alignment
        for text, bold, italic in runs:
            run = paragraph.add_run()
            run.text = text
            run.font.bold = bold
            run.font.italic = italic
            run.font.size = Pt(_clamp(font_size, 8.0, 96.0))
            if color is not None:
                run.font.color.rgb = color


async def _render_image_block(slide, block: Dict[str, Any], left: Emu, top: Emu, width: Emu, height: Emu) -> None:
    image_path = await _materialize_image_source(block.get("prompt"))
    if not image_path:
        logger.warning("Skipping image block '%s': image path not found (%s)", block.get("id"), block.get("prompt"))
        return

    try:
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    except Exception as exc:  # noqa: BLE001
        logger.warning("Skipping image block '%s': failed to add image (%s)", block.get("id"), exc)


def _chart_type_to_pptx(chart_type: str) -> XL_CHART_TYPE:
    normalized = (chart_type or "bar").lower()
    if normalized == "line":
        return XL_CHART_TYPE.LINE
    if normalized == "pie":
        return XL_CHART_TYPE.PIE
    if normalized == "donut":
        return XL_CHART_TYPE.DOUGHNUT
    if normalized == "area":
        return XL_CHART_TYPE.AREA
    return XL_CHART_TYPE.COLUMN_CLUSTERED


def _render_chart_block(slide, block: Dict[str, Any], left: Emu, top: Emu, width: Emu, height: Emu) -> None:
    points = block.get("data") if isinstance(block.get("data"), list) else []
    if not points:
        logger.warning("Skipping chart block '%s': missing data", block.get("id"))
        return

    labels: List[str] = []
    values: List[float] = []
    for point in points:
        if not isinstance(point, dict):
            continue
        labels.append(str(point.get("label", "")))
        values.append(_as_float(point.get("value"), 0.0))

    if not labels:
        logger.warning("Skipping chart block '%s': invalid chart data points", block.get("id"))
        return

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(str(block.get("title") or "Series"), values)

    try:
        xl_type = _chart_type_to_pptx(str(block.get("chart_type") or "bar"))
        graphic = slide.shapes.add_chart(
            xl_type,
            left,
            top,
            width,
            height,
            chart_data,
        )
        chart = graphic.chart
        # Pie/doughnut + legend is a common source of strict OOXML issues in viewers.
        chart.has_legend = xl_type not in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)
    except Exception as exc:  # noqa: BLE001
        logger.warning("Skipping chart block '%s': failed to add chart (%s)", block.get("id"), exc)


async def export_spatial_to_pptx(
    *,
    editor_payload: Dict[str, Any],
    output_path: str,
) -> str:
    """Export spatial JSON editor payload to a PPTX file."""

    slides = editor_payload.get("slides") if isinstance(editor_payload, dict) else None
    if not isinstance(slides, list) or not slides:
        raise ValueError("editor_payload.slides is required and must be a non-empty list")

    prs = Presentation()
    blank_layout = _blank_slide_layout(prs)

    for slide_payload in slides:
        slide = prs.slides.add_slide(blank_layout)
        blocks = slide_payload.get("blocks") if isinstance(slide_payload, dict) else []
        if not isinstance(blocks, list):
            blocks = []

        # Draw lower z-index first so higher z-index appears on top.
        sorted_blocks = sorted(
            blocks,
            key=lambda b: _as_float((b or {}).get("z_index"), 0.0) if isinstance(b, dict) else 0.0,
        )

        for block in sorted_blocks:
            if not isinstance(block, dict):
                continue

            x, y, w, h = _normalize_position(block.get("position"))
            left, top, width, height = _percent_to_emu(prs, x, y, w, h)
            block_type = str(block.get("type") or "").lower()

            if block_type == "text":
                _render_text_block(slide, block, left, top, width, height)
            elif block_type == "image":
                await _render_image_block(slide, block, left, top, width, height)
            elif block_type == "chart":
                _render_chart_block(slide, block, left, top, width, height)
            else:
                logger.warning("Skipping unknown block type '%s' on slide", block_type)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
